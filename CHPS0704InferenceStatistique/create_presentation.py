"""
Script de création de la présentation PowerPoint
Projet : Inférence Statistique - Analyse de données cardiaques
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Couleurs du thème
BLEU_FONCE = RgbColor(0x1E, 0x3A, 0x5F)  # Titres
BLEU_CLAIR = RgbColor(0x3D, 0x5A, 0x80)  # Sous-titres
GRIS = RgbColor(0x6B, 0x70, 0x7B)         # Texte secondaire
BLANC = RgbColor(0xFF, 0xFF, 0xFF)

def set_slide_background(slide, color=RgbColor(0xF8, 0xF9, 0xFA)):
    """Définit un fond légèrement gris"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_shape(slide, text, top=0.3, font_size=36, color=BLEU_FONCE):
    """Ajoute un titre centré"""
    left = Inches(0.5)
    top = Inches(top)
    width = Inches(9)
    height = Inches(0.8)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_subtitle(slide, text, top=1.0, font_size=18, color=GRIS):
    """Ajoute un sous-titre"""
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.5))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_bullet_points(slide, points, left=0.5, top=1.8, width=9, font_size=20):
    """Ajoute une liste de points"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4))
    tf = shape.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RgbColor(0x33, 0x33, 0x33)
        p.space_after = Pt(12)
    return shape

def add_image_centered(slide, image_path, top=1.5, max_width=8.5, max_height=5):
    """Ajoute une image centrée"""
    if not os.path.exists(image_path):
        print(f"Image non trouvée: {image_path}")
        return None

    # Calculer la taille pour garder les proportions
    from PIL import Image
    img = Image.open(image_path)
    img_width, img_height = img.size
    aspect = img_width / img_height

    if aspect > max_width / max_height:
        width = max_width
        height = width / aspect
    else:
        height = max_height
        width = height * aspect

    left = (10 - width) / 2  # Centrer horizontalement
    slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))

def add_table(slide, data, left=1, top=1.8, col_widths=None):
    """Ajoute un tableau stylisé"""
    rows = len(data)
    cols = len(data[0])

    width = Inches(8) if col_widths is None else Inches(sum(col_widths))
    height = Inches(rows * 0.5)

    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), width, height).table

    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)

            # Style
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14)
            para.alignment = PP_ALIGN.CENTER

            # En-tête en bleu
            if i == 0:
                para.font.bold = True
                para.font.color.rgb = BLANC
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLEU_FONCE
            else:
                para.font.color.rgb = RgbColor(0x33, 0x33, 0x33)

    return table

def create_presentation():
    """Crée la présentation complète"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Chemin des images
    fig_path = "figures"

    # ============================================
    # SLIDE 1: Titre
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_background(slide, RgbColor(0x1E, 0x3A, 0x5F))

    add_title_shape(slide, "Inférence Statistique", top=2.0, font_size=44, color=BLANC)
    add_subtitle(slide, "Analyse et Modélisation de Données Cardiaques", top=2.8, font_size=24, color=RgbColor(0xA0, 0xC4, 0xE8))

    # Auteur
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Nicolas Marano"
    p.font.size = Pt(20)
    p.font.color.rgb = BLANC
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Master 1 - CHPS"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RgbColor(0xA0, 0xC4, 0xE8)
    p2.alignment = PP_ALIGN.CENTER

    # ============================================
    # SLIDE 2: Contexte
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Contexte du Projet", top=0.3)

    points = [
        "Analyse statistique de données de santé cardiaque",
        "Dataset : 500 patients avec mesures médicales",
        "Objectif : Appliquer les méthodes d'inférence statistique",
        "Données synthétiques avec relations médicales réalistes"
    ]
    add_bullet_points(slide, points, font_size=22)

    # ============================================
    # SLIDE 3: Dataset
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Description du Dataset", top=0.3)

    data = [
        ["Variable", "Type", "Unité", "Moyenne", "Variance"],
        ["Âge", "Quantitative", "années", "54.57", "133.38"],
        ["Cholestérol", "Quantitative", "mg/dL", "211.37", "1510.64"],
        ["Fréq. cardiaque", "Quantitative", "bpm", "75.87", "141.85"],
        ["Pression artérielle", "Quantitative", "mmHg", "133.48", "193.16"],
        ["Maladie cardiaque", "Qualitative", "0/1", "39.8%", "-"]
    ]
    add_table(slide, data, left=1, top=1.5)

    # ============================================
    # SLIDE 4: Distributions
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Analyse Descriptive", top=0.2)
    add_subtitle(slide, "Distribution des variables quantitatives", top=0.9, font_size=16)
    add_image_centered(slide, f"{fig_path}/histogrammes.png", top=1.3, max_height=5.5)

    # ============================================
    # SLIDE 5: Ajustement Normal
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Test de Normalité", top=0.2)
    add_subtitle(slide, "Ajustement à la loi normale - Cholestérol", top=0.9, font_size=16)
    add_image_centered(slide, f"{fig_path}/ajustement_cholesterol.png", top=1.3, max_height=5.5)

    # ============================================
    # SLIDE 6: Estimation
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Estimation des Paramètres", top=0.3)

    # Sous-titre méthodes
    add_subtitle(slide, "Comparaison : Méthode des Moments vs Maximum de Vraisemblance", top=1.0, font_size=16)

    data = [
        ["Méthode", "μ (mg/dL)", "σ² (mg/dL)²"],
        ["Moments (MM)", "211.37", "1510.64"],
        ["Vraisemblance (MLE)", "211.37", "1507.62"]
    ]
    add_table(slide, data, left=2, top=1.8)

    points = [
        "Les deux méthodes donnent des résultats quasi-identiques",
        "Biais MLE négligeable : -0.047",
        "Test Kolmogorov-Smirnov : hypothèse normale acceptée"
    ]
    add_bullet_points(slide, points, top=4.0, font_size=18)

    # ============================================
    # SLIDE 7: Intervalles de Confiance
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Intervalles de Confiance (95%)", top=0.2)
    add_image_centered(slide, f"{fig_path}/intervalles_confiance.png", top=1.2, max_height=5.5)

    # ============================================
    # SLIDE 8: Tests d'Hypothèses
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Tests d'Hypothèses", top=0.2)
    add_subtitle(slide, "Comparaison entre groupes : malades vs sains", top=0.9, font_size=16)
    add_image_centered(slide, f"{fig_path}/tests_hypotheses.png", top=1.3, max_height=5.5)

    # ============================================
    # SLIDE 9: Scatter Plot
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Régression Linéaire", top=0.2)
    add_subtitle(slide, "Relation : Âge → Pression Artérielle", top=0.9, font_size=16)
    add_image_centered(slide, f"{fig_path}/scatter_plot.png", top=1.3, max_height=5.5)

    # ============================================
    # SLIDE 10: Régression
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Modèle de Régression", top=0.2)
    add_image_centered(slide, f"{fig_path}/regression_lineaire.png", top=1.0, max_height=4.5)

    # Équation
    shape = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(0.5))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Pression = 86.81 + 0.855 × Âge    |    R² = 75%"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLEU_FONCE
    p.alignment = PP_ALIGN.CENTER

    # ============================================
    # SLIDE 11: Résidus
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Validation du Modèle", top=0.2)
    add_subtitle(slide, "Analyse des résidus", top=0.9, font_size=16)
    add_image_centered(slide, f"{fig_path}/analyse_residus.png", top=1.3, max_height=5.5)

    # ============================================
    # SLIDE 12: Résultats Clés
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title_shape(slide, "Résultats Clés", top=0.3)

    points = [
        "Cholestérol significativement > 200 mg/dL (p < 0.001)",
        "Différence significative malades vs sains (p < 0.01)",
        "Pression artérielle augmente de ~0.86 mmHg/an",
        "75% de la variance expliquée par l'âge",
        "Hypothèses du modèle validées"
    ]
    add_bullet_points(slide, points, font_size=24, top=1.5)

    # ============================================
    # SLIDE 13: Conclusion
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, RgbColor(0x1E, 0x3A, 0x5F))

    add_title_shape(slide, "Conclusion", top=1.5, font_size=40, color=BLANC)

    shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(3))
    tf = shape.text_frame
    tf.word_wrap = True

    conclusions = [
        "Application complète des méthodes d'inférence",
        "Résultats cohérents avec les attentes médicales",
        "Modèle prédictif fiable pour la pression artérielle",
        "Limites : données synthétiques, corrélation ≠ causalité"
    ]

    for i, text in enumerate(conclusions):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {text}"
        p.font.size = Pt(22)
        p.font.color.rgb = BLANC
        p.space_after = Pt(16)
        p.alignment = PP_ALIGN.CENTER

    # Sauvegarder
    output_path = "presentation_inference_statistique.pptx"
    prs.save(output_path)
    print(f"Présentation créée avec succès : {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
